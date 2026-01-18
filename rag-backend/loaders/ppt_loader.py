from pptx import Presentation
from loaders.base_loader import BaseLoader

class PPTLoader(BaseLoader):

    def load(self, file_path: str) -> str:
        prs = Presentation(file_path)
        slides_text = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slides_text.append(shape.text)

        return "\n".join(slides_text)
